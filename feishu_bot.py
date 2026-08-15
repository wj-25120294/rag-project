# -*- coding: utf-8 -*-
"""
飞书群聊 @机器人 生成 PPT —— 独立机器人脚本
================================================
功能:在飞书群聊里 @机器人 提问(如"用langchain学习路线做一份PPT"),
机器人调用 DeepSeek 生成 PPT 大纲 → 用 python-pptx 生成 .pptx 文件
→ 上传到飞书 → 作为文件消息发回群里。

本文件完全独立,不依赖本项目其它代码(RAG / milvus / fastapi 等)。

【一】安装依赖(与本文件同目录执行)
    pip install lark-oapi python-pptx requests python-dotenv

【二】飞书开放平台配置(open.feishu.cn → 开发者后台 → 你的自建应用)
    1. 权限管理,添加:
       im:message             接收群聊中@机器人消息
       im:message:send_as_bot 以应用身份发消息
       im:resource            上传文件到飞书(新版可能叫 im:file / im:resource,按提示勾选)
    2. 事件订阅:订阅方式选「使用长连接接收事件」,添加事件 im.message.receive_v1(接收消息)
       —— 选长连接,就不需要公网服务器 / HTTPS 回调地址,本地跑即可
    3. 发布应用版本,并把机器人拉进群(群设置 → 群机器人 → 添加机器人)
    4. 在本文件同目录的 .env 里填:
       FEISHU_APP_ID=...
       FEISHU_APP_SECRET=...
       DEEPSEEK_API_KEY=...

【三】运行
    python feishu_bot.py

【四】使用
    群里 @机器人 提问即可,机器人先回复"收到,正在生成...",完成后把 PPT 文件发到群里。
"""

import json
import os
import re
import threading
import time

import requests
from dotenv import load_dotenv

# ---------------------------------------------------------------- 配置
load_dotenv()  # 读取同目录 .env(FEISHU_APP_ID / FEISHU_APP_SECRET / DEEPSEEK_API_KEY)

APP_ID = os.getenv("FEISHU_APP_ID")
APP_SECRET = os.getenv("FEISHU_APP_SECRET")
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")

FEISHU_OPEN_URL = "https://open.feishu.cn/open-apis"
DEEPSEEK_URL = "https://api.deepseek.com/chat/completions"

MAX_SLIDES = 8   # PPT 最多页数(不含封面)
PPT_PROMPT = """你是PPT大纲设计专家。根据用户主题,输出一份结构清晰、内容专业的PPT大纲。
要求:
- 只输出一个 JSON 对象,不要输出 JSON 以外的任何文字
- JSON 格式如下:
{{"title": "演示标题", "slides": [{{"title": "页面标题", "bullets": ["要点1", "要点2", "要点3"]}}]}}
- slides 不超过 {max_slides} 页,每页 bullets 3~5 条,每条不超过30字
- 主题: {topic}"""

# ---------------------------------------------------------------- 运行日志(写入 bot.log,便于排查)
LOG_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bot.log")


def _log(msg):
    print(msg, flush=True)
    try:
        with open(LOG_PATH, "a", encoding="utf-8") as f:
            f.write(time.strftime("%Y-%m-%d %H:%M:%S ") + msg + "\n")
    except Exception:
        pass


def _log_exc():
    import traceback
    _log("[EXCEPTION]\n" + traceback.format_exc())

# ---------------------------------------------------------------- 飞书 token(带缓存)
_token_cache = {"token": None, "expire_at": 0}
_bot_open_id = None


def get_tenant_access_token():
    """获取(并缓存)tenant_access_token,用于所有飞书 API 调用"""
    now = time.time()
    if _token_cache["token"] and now < _token_cache["expire_at"]:
        return _token_cache["token"]
    def _do():
        resp = requests.post(
            f"{FEISHU_OPEN_URL}/auth/v3/tenant_access_token/internal",
            json={"app_id": APP_ID, "app_secret": APP_SECRET}, timeout=15,
        )
        resp.raise_for_status()
        return resp.json()
    data = _call_with_retry(_do)
    if data.get("code") != 0:
        raise RuntimeError(f"获取 tenant_access_token 失败: {data}")
    _token_cache["token"] = data["tenant_access_token"]
    _token_cache["expire_at"] = now + data["expire"] - 60  # 提前 60 秒过期
    return _token_cache["token"]


def _headers():
    return {"Authorization": f"Bearer {get_tenant_access_token()}"}


def _call_with_retry(func, *args, retries=4, **kwargs):
    """调用 func;网络异常(SSL瞬断)、空响应/非JSON、4xx/5xx 都自动重试;
    业务层错误(如权限 99991672)不重试。调用方需在 func 内 raise_for_status + 解析JSON。"""
    last_exc = None
    for attempt in range(retries):
        try:
            return func(*args, **kwargs)
        except (requests.exceptions.RequestException, ValueError) as e:
            last_exc = e
            _log(f"[网络重试] {attempt + 1}/{retries} 失败: {e},2秒后重试")
            time.sleep(2)
    raise last_exc


def get_bot_open_id():
    """获取机器人自己的 open_id,用来判断群里是否 @ 了机器人"""
    global _bot_open_id
    if _bot_open_id:
        return _bot_open_id
    try:
        resp = requests.get(f"{FEISHU_OPEN_URL}/bot/v3/info", headers=_headers(), timeout=10)
        data = resp.json()
        if data.get("code") == 0:
            _bot_open_id = data["bot"]["open_id"]
    except Exception:
        pass  # 取不到也没关系,群里只要出现 @ 就当作在叫机器人
    return _bot_open_id


# ---------------------------------------------------------------- 飞书发送 / 上传
def reply_text(message_id, text):
    """在消息下回复一条文本(会挂到原消息下面)"""
    payload = {"msg_type": "text", "content": json.dumps({"text": text})}
    def _do():
        resp = requests.post(
            f"{FEISHU_OPEN_URL}/im/v1/messages/{message_id}/reply",
            headers={**_headers(), "Content-Type": "application/json"},
            json=payload, timeout=30,
        )
        resp.raise_for_status()
        return resp.json()
    return _call_with_retry(_do)


def upload_file(file_path, file_name):
    """把本地文件上传到飞书,返回 file_key(失败自动重试,每次重新打开文件)"""
    def _do_upload():
        with open(file_path, "rb") as f:
            resp = requests.post(
                f"{FEISHU_OPEN_URL}/im/v1/files",
                headers=_headers(),
                data={"file_type": "stream", "file_name": file_name},
                files={"file": f},
                timeout=60,
            )
        resp.raise_for_status()
        return resp.json()
    data = _call_with_retry(_do_upload)
    if data.get("code") != 0:
        raise RuntimeError(f"上传文件失败: {data}")
    return data["data"]["file_key"]


def send_file(chat_id, file_key):
    """往群聊发送一个文件消息"""
    content = json.dumps({"file_key": file_key})
    def _do():
        resp = requests.post(
            f"{FEISHU_OPEN_URL}/im/v1/messages?receive_id_type=chat_id",
            headers={**_headers(), "Content-Type": "application/json"},
            json={"receive_id": chat_id, "msg_type": "file", "content": content},
            timeout=30,
        )
        resp.raise_for_status()
        return resp.json()
    data = _call_with_retry(_do)
    if data.get("code") != 0:
        raise RuntimeError(f"发送文件消息失败: {data}")


# ---------------------------------------------------------------- DeepSeek 生成大纲
def generate_outline(topic):
    """调用 DeepSeek 生成 PPT 大纲,返回 dict: {"title":..., "slides":[...]}"""
    messages = [
        {"role": "system", "content": "你只输出 JSON,不要输出任何其他内容。"},
        {"role": "user", "content": PPT_PROMPT.format(topic=topic, max_slides=MAX_SLIDES)},
    ]
    def _do():
        resp = requests.post(
            DEEPSEEK_URL,
            headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"},
            json={
                "model": "deepseek-chat",
                "messages": messages,
                "temperature": 0.7,
                "response_format": {"type": "json_object"},
            },
            timeout=180,
        )
        resp.raise_for_status()
        return resp.json()
    data = _call_with_retry(_do, retries=3)
    if "choices" not in data:
        raise RuntimeError(f"DeepSeek 调用失败: {data}")

    text = data["choices"][0]["message"]["content"]
    # 兼容模型偶尔输出 ```json ... ``` 包裹的情况
    text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text.strip())
    return json.loads(text)


# ---------------------------------------------------------------- python-pptx 生成文件
def build_pptx(outline, out_path):
    """按大纲生成 .pptx 文件(封面 + 内容页)"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title = outline.get("title") or "演示文稿"
    slides = outline.get("slides") or []

    # 封面页
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = "By Feishu Bot"

    # 内容页
    for idx, item in enumerate(slides, 1):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"{idx}. {item.get('title', '')}"
        body = s.placeholders[1].text_frame
        body.clear()
        for bullet in item.get("bullets", []):
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    prs.save(out_path)


# ---------------------------------------------------------------- 消息处理
def extract_text(message):
    """从消息里提取纯文本提问,兼容 text / post 等消息类型和不同结构"""
    msg_type = getattr(message, "message_type", "") or getattr(message, "msg_type", "")
    content = getattr(message, "content", "")

    # content 可能是 JSON 字符串,也可能是 SDK 已解析好的 dict
    if isinstance(content, str):
        try:
            body = json.loads(content)
        except Exception:
            body = {}
    elif isinstance(content, dict):
        body = content
    else:
        body = {}

    text = ""
    if isinstance(body.get("text"), str):
        text = body["text"]
    else:
        # post 消息:content 是 [[{tag,text}, ...], ...] 的结构
        lines = body.get("content")
        if isinstance(lines, list):
            parts = []
            for line in lines:
                if not isinstance(line, list):
                    continue
                for seg in line:
                    if isinstance(seg, dict) and seg.get("tag") == "text" and isinstance(seg.get("text"), str):
                        parts.append(seg["text"])
            text = "".join(parts)

    # 去掉 "@_user_xxx" 这类 @ 占位符后,剩下的才是问题
    text = re.sub(r"@_user_\w+", "", text).strip()
    _log(f"[DEBUG] msg_type={msg_type} | 解析出的问题=[{text}]")
    return text


def is_mention_bot(message, chat_type):
    """是否处理这条消息。
    飞书默认只把「@机器人」的群消息推给应用(不@根本收不到),
    所以这里直接放行,不再依赖 @ 识别 —— 最简单也最稳。
    mentions 仅在日志里用于排查。"""
    return True


def build_and_send(chat_id, message_id, question):
    """生成 PPT 并作为文件发到群聊(在后台线程执行)"""
    try:
        outline = generate_outline(question)
        out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ppt_output")
        os.makedirs(out_dir, exist_ok=True)
        fname = f"ppt_{int(time.time())}.pptx"
        out_path = os.path.join(out_dir, fname)
        build_pptx(outline, out_path)

        file_key = upload_file(out_path, fname)
        send_file(chat_id, file_key)
    except Exception as e:
        _log_exc()
        try:
            reply_text(message_id, f"生成 PPT 失败:{e}")
        except Exception:
            pass


def process_message(data):
    """收到消息事件后的处理入口"""
    try:
        event = getattr(data, "event", None)
        message = getattr(event, "message", None)
        sender = getattr(event, "sender", None)
        if not message or not message.content:
            _log("[DEBUG] 收到空消息,忽略")
            return

        # 忽略机器人自己发的消息
        if getattr(sender, "sender_type", "") == "bot":
            _log("[DEBUG] 忽略机器人自己的消息")
            return

        message_id = message.message_id
        chat_id = message.chat_id
        chat_type = getattr(message, "chat_type", "")
        mentions = [
            f"{getattr(m, 'key', '')}/{getattr(getattr(m, 'id', None), 'open_id', '')}"
            for m in (getattr(message, "mentions", None) or [])
        ]
        _log(f"[DEBUG] 收到消息 | id={message_id} | chat_type={chat_type} | msg_type={message.message_type} | mentions={mentions} | 内容={message.content[:80]}")

        if not is_mention_bot(message, chat_type):
            _log("[DEBUG] 未被@或非目标消息,跳过")
            return

        question = extract_text(message)
        if not question:
            reply_text(message_id, "请告诉我想做什么主题的PPT,例如:@机器人 用langchain学习路线做一份PPT")
            return

        # 先快速回复,再放到后台线程慢慢生成
        reply_text(message_id, "收到,正在生成 PPT,请稍等几秒~")
        threading.Thread(target=build_and_send, args=(chat_id, message_id, question), daemon=True).start()
    except Exception:
        _log_exc()


# ---------------------------------------------------------------- 长连接启动
def on_message_receive(data) -> None:
    """lark-oapi 事件回调(im.message.receive_v1)"""
    process_message(data)


def main():
    if not all([APP_ID, APP_SECRET, DEEPSEEK_API_KEY]):
        print("缺少配置:请在同目录 .env 中设置 FEISHU_APP_ID / FEISHU_APP_SECRET / DEEPSEEK_API_KEY")
        return

    import lark_oapi as lark

    event_handler = (
        lark.EventDispatcherHandler.builder("", "")   # 长连接模式不需要 token/encrypt_key
        .register_p2_im_message_receive_v1(on_message_receive)
        .build()
    )
    ws_client = lark.ws.Client(APP_ID, APP_SECRET, event_handler=event_handler,
                               log_level=lark.LogLevel.INFO)
    _log("机器人已启动,正在连接飞书长连接...(Ctrl+C 退出)")
    ws_client.start()


if __name__ == "__main__":
    main()
